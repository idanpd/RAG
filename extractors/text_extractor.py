"""
Text content extractor for documents, PDFs, spreadsheets, etc.
"""

from pathlib import Path
from typing import Union, Dict, Any
import logging

from .base import BaseExtractor, ExtractionResult

# Import libraries with fallbacks
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)


class TextExtractor(BaseExtractor):
    """Extractor for text-based documents."""
    
    SUPPORTED_EXTENSIONS = {
        '.txt', '.md', '.csv', '.json', '.log', '.py', '.js', '.html', '.xml',
        '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', 
        '.odt', '.ods', '.odp', '.rtf'
    }
    
    def __init__(self, config: Dict[str, Any] = None):
        super().__init__(config)
        self.max_file_size = self.config.get('max_text_file_size', 100 * 1024 * 1024)  # 100MB
    
    def can_extract(self, file_path: Union[str, Path]) -> bool:
        """Check if this extractor can handle the file."""
        path = Path(file_path)
        return path.suffix.lower() in self.SUPPORTED_EXTENSIONS
    
    def extract(self, file_path: Union[str, Path]) -> ExtractionResult:
        """Extract text content from various document formats."""
        path = Path(file_path)
        
        # Check file size
        try:
            if path.stat().st_size > self.max_file_size:
                return ExtractionResult(
                    content="",
                    success=False,
                    error_message=f"File too large: {path.stat().st_size} bytes"
                )
        except Exception as e:
            return ExtractionResult(
                content="",
                success=False,
                error_message=f"Cannot access file: {e}"
            )
        
        suffix = path.suffix.lower()
        
        # Route to appropriate extraction method
        if suffix == '.pdf':
            return self._extract_pdf(path)
        elif suffix in {'.docx', '.doc'}:
            return self._extract_docx(path)
        elif suffix in {'.xlsx', '.xls', '.csv'}:
            return self._extract_spreadsheet(path)
        elif suffix in {'.pptx', '.ppt'}:
            return self._extract_presentation(path)
        elif suffix in {'.txt', '.md', '.py', '.js', '.html', '.xml', '.json', '.log'}:
            return self._extract_text_file(path)
        else:
            # Try as text file first, then fallback
            result = self._extract_text_file(path)
            if result.success and result.content.strip():
                return result
            return self._create_fallback_result(path)
    
    def _extract_pdf(self, path: Path) -> ExtractionResult:
        """Extract text from PDF files."""
        if not PYMUPDF_AVAILABLE:
            return ExtractionResult(
                content="",
                success=False,
                error_message="PyMuPDF not available for PDF extraction"
            )
        
        try:
            doc = fitz.open(str(path))
            pages = []
            metadata = {
                'page_count': doc.page_count,
                'title': doc.metadata.get('title', ''),
                'author': doc.metadata.get('author', ''),
                'subject': doc.metadata.get('subject', '')
            }
            
            for page_num in range(doc.page_count):
                page = doc[page_num]
                text = page.get_text()
                if text.strip():
                    pages.append(f"--- Page {page_num + 1} ---\n{text}")
            
            doc.close()
            
            content = "\n\n".join(pages)
            return ExtractionResult(
                content=content,
                metadata=metadata,
                success=True,
                extraction_method='pymupdf'
            )
            
        except Exception as e:
            return ExtractionResult(
                content="",
                success=False,
                error_message=f"PDF extraction failed: {e}",
                extraction_method='pymupdf'
            )
    
    def _extract_docx(self, path: Path) -> ExtractionResult:
        """Extract text from DOCX files."""
        if not DOCX_AVAILABLE:
            return ExtractionResult(
                content="",
                success=False,
                error_message="python-docx not available for DOCX extraction"
            )
        
        try:
            doc = docx.Document(str(path))
            paragraphs = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(" | ".join(row_data))
                if table_data:
                    tables.append("\n".join(table_data))
            
            content_parts = []
            if paragraphs:
                content_parts.append("\n\n".join(paragraphs))
            if tables:
                content_parts.append("\n\n=== TABLES ===\n\n" + "\n\n".join(tables))
            
            content = "\n\n".join(content_parts)
            
            return ExtractionResult(
                content=content,
                metadata={
                    'paragraph_count': len(paragraphs),
                    'table_count': len(tables)
                },
                success=True,
                extraction_method='python-docx'
            )
            
        except Exception as e:
            return ExtractionResult(
                content="",
                success=False,
                error_message=f"DOCX extraction failed: {e}",
                extraction_method='python-docx'
            )
    
    def _extract_spreadsheet(self, path: Path) -> ExtractionResult:
        """Extract data from spreadsheet files."""
        if not PANDAS_AVAILABLE:
            return ExtractionResult(
                content="",
                success=False,
                error_message="pandas not available for spreadsheet extraction"
            )
        
        try:
            suffix = path.suffix.lower()
            
            if suffix == '.csv':
                df = pd.read_csv(str(path))
                sheets = {'Sheet1': df}
            else:
                sheets = pd.read_excel(str(path), sheet_name=None)
            
            content_parts = []
            for sheet_name, df in sheets.items():
                if df.empty:
                    continue
                
                content_parts.append(f"=== {sheet_name} ===")
                
                # Add column headers
                headers = " | ".join(df.columns.astype(str))
                content_parts.append(headers)
                content_parts.append("-" * len(headers))
                
                # Add data rows (limit to prevent huge outputs)
                max_rows = self.config.get('max_spreadsheet_rows', 1000)
                for idx, row in df.head(max_rows).iterrows():
                    row_text = " | ".join(row.astype(str).fillna(""))
                    content_parts.append(row_text)
                
                if len(df) > max_rows:
                    content_parts.append(f"... ({len(df) - max_rows} more rows)")
            
            content = "\n".join(content_parts)
            
            return ExtractionResult(
                content=content,
                metadata={
                    'sheet_count': len(sheets),
                    'total_rows': sum(len(df) for df in sheets.values())
                },
                success=True,
                extraction_method='pandas'
            )
            
        except Exception as e:
            return ExtractionResult(
                content="",
                success=False,
                error_message=f"Spreadsheet extraction failed: {e}",
                extraction_method='pandas'
            )
    
    def _extract_presentation(self, path: Path) -> ExtractionResult:
        """Extract text from presentation files."""
        if not PPTX_AVAILABLE:
            return ExtractionResult(
                content="",
                success=False,
                error_message="python-pptx not available for presentation extraction"
            )
        
        try:
            prs = Presentation(str(path))
            slides = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                slide_text.append(f"--- Slide {i + 1} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # More than just the slide header
                    slides.append("\n".join(slide_text))
            
            content = "\n\n".join(slides)
            
            return ExtractionResult(
                content=content,
                metadata={'slide_count': len(prs.slides)},
                success=True,
                extraction_method='python-pptx'
            )
            
        except Exception as e:
            return ExtractionResult(
                content="",
                success=False,
                error_message=f"Presentation extraction failed: {e}",
                extraction_method='python-pptx'
            )
    
    def _extract_text_file(self, path: Path) -> ExtractionResult:
        """Extract content from plain text files."""
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(path, 'r', encoding=encoding, errors='ignore') as f:
                    content = f.read()
                
                return ExtractionResult(
                    content=content,
                    metadata={'encoding': encoding, 'size': len(content)},
                    success=True,
                    extraction_method='text_file'
                )
                
            except Exception as e:
                continue
        
        return ExtractionResult(
            content="",
            success=False,
            error_message="Failed to read text file with any encoding",
            extraction_method='text_file'
        )